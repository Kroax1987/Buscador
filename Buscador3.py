import streamlit as st
import pandas as pd
import os
from pptx import Presentation

def search_in_excel(file, keyword):
    try:
        df = pd.read_excel(file, sheet_name=None)
    except Exception as e:
        st.error(f"Erro lendo Excel: {e}")
        return []
    results = []
    for sheet, data in df.items():
        mask = data.applymap(lambda x: keyword.lower() in str(x).lower() if pd.notnull(x) else False)
        matched = data[mask.any(axis=1)]
        if not matched.empty:
            results.append((sheet, matched))
    return results

def search_in_txt(file, keyword):
    try:
        text = file.read().decode('utf-8')
    except Exception as e:
        st.error(f"Erro lendo TXT: {e}")
        return []
    lines = text.splitlines()
    results = [line for line in lines if keyword.lower() in line.lower()]
    return results

def search_in_pptx(file, keyword):
    try:
        prs = Presentation(file)
    except Exception as e:
        st.error(f"Erro lendo PPTX: {e}")
        return []
    results = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and keyword.lower() in shape.text.lower():
                results.append(f"Slide {i+1}: {shape.text}")
    return results

st.title("Buscador Inteligente")

uploaded_files = st.file_uploader("Faça upload dos arquivos", accept_multiple_files=True)
keyword = st.text_input("Digite a palavra-chave para busca")

if uploaded_files and keyword:
    for uploaded_file in uploaded_files:
        st.subheader(f"Resultados no arquivo: {uploaded_file.name}")
        ext = os.path.splitext(uploaded_file.name)[1].lower()
        
        if ext in ['.xls', '.xlsx']:
            results = search_in_excel(uploaded_file, keyword)
            if results:
                for sheet, df in results:
                    st.write(f"Na aba {sheet}:")
                    st.dataframe(df)
            else:
                st.write("Nenhum resultado encontrado.")
        
        elif ext == '.txt':
            results = search_in_txt(uploaded_file, keyword)
            if results:
                for line in results:
                    st.write(line)
            else:
                st.write("Nenhum resultado encontrado.")
        
        elif ext == '.pptx':
            results = search_in_pptx(uploaded_file, keyword)
            if results:
                for res in results:
                    st.write(res)
            else:
                st.write("Nenhum resultado encontrado.")
        
        else:
            st.write("Formato de arquivo não suportado.")
